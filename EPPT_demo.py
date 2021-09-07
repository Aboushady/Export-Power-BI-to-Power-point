import sys
import ssas_api as powerbi
import pandas as pd


"""
Part 0:
    Refer to this guid for info regarding the connection to Power BI's data base.
    https://dataveld.com/2020/07/20/python-as-an-external-tool-for-power-bi-desktop-part-1/
"""

print('Power BI Desktop Connection')
print(str(sys.argv[1]))
print(str(sys.argv[2]))

conn = "\nProvider=MSOLAP;Data Source=" + str(sys.argv[1]) + ";Initial Catalog='';"
print(conn)


dax_string = 'EVALUATE ROW("Loading .NET assemblies",1)'
df = powerbi.get_DAX(connection_string=conn, dax_string=dax_string)

print("\nCrossing the streams...")
global System, DataTable, AMO, ADOMD


"""
Part 1:
    Get the databases info, and retrieve the Meta data for the indended one,
    using the Microsoft Analysis services DLLs.
"""


import System
from System.Data import DataTable
import Microsoft.AnalysisServices.Tabular as TOM
import Microsoft.AnalysisServices.AdomdClient as ADOMD
from pptx import Presentation
from pptx.util import Inches
import plotly.express as px
from pandas.plotting import table
import imgkit

print("\nReticulating splines...")


# Connecting to Tabular AnalysisServices.
TOMServer = TOM.Server()
TOMServer.Connect(conn)

print("\nConnectoin Successfully to TOM's Server...")
print()

# Show Databases info
for item in TOMServer.Databases:
    print("Database: ", item.Name)
    print("Compatibility Level: ", item.CompatibilityLevel)
    print("Created: ", item.CreatedTimestamp)

input("Click to cont.")

'''
Retrieve the Meta data for the database we got from the connection earlier,
from TOM's databases.
'''

DatabaseId = str(sys.argv[2])
PowerBIDatabase = TOMServer.Databases[DatabaseId]

print()

'''
The following is just to take a look at the retrieved Meta data.
'''

# Define measure dataframe
dfMeasures = pd.DataFrame(
    columns=['Table',
             'Name',
             'Description',
             'DataType',
             'DataCategory',
             'Expression',
             'FormatString',
             'DisplayFolder',
             'Implicit',
             'Hidden',
             'ModifiedTime',
             'State'])

# Define column dataframe
dfColumns = pd.DataFrame(
    columns=['Table',
             'Name'])

# Tables
print("Listing tables...")
for table in PowerBIDatabase.Model.Tables:
    print(table.Name + "\n")

    # Assign current table by name
    CurrentTable = PowerBIDatabase.Model.Tables.Find(table.Name)


    # Measures
    # Get the measures from each table and append it to the measures dataframe.
    for measure in CurrentTable.Measures:
        new_row = {'Table':table.Name,
                'Name':measure.Name,
                'Description':measure.Description,
                'DataType':measure.DataType,
                'DataCategory':measure.DataCategory,
                'Expression':measure.Expression,
                'FormatString':measure.FormatString,
                'DisplayFolder':measure.DisplayFolder,
                'Implicit':measure.IsSimpleMeasure,
                'Hidden':measure.IsHidden,
                'ModifiedTime':measure.ModifiedTime,
                'State':measure.State}
        dfMeasures = dfMeasures.append(new_row, ignore_index=True)

    # Columns
    # Get the columns from each table and append it to the columns dataframe.
    for column in CurrentTable.Columns:
        new_row = {'Table':table.Name,
                'Name':column.Name}

        dfColumns = dfColumns.append(new_row, ignore_index=True)

print(dfMeasures)
print(dfColumns)


# A dax query that returns a table, here grouped by the "Year".
# Sheet1 : is the name of the table in Power BI's data model.
dax_q = """
    Evaluate 
    SUMMARIZECOLUMNS(Sheet1[Year],
                        "Total Purchases", SUMX(Sheet1, Sheet1[Number of Purchases]))

"""

###################################################################################

'''
Execute the dax query, and convert the table into a data frame.
# connection_string : is the connection string to Power BI's Meta data.
# dax_string : the dax query.
# dax_exec_q : Holds the returned data frame.
'''
dax_exec_q = powerbi.get_DAX(connection_string=conn, dax_string=dax_q)

###################################################################################

'''
Create a plotly pie chart figure to visualize the data frame.
'''
# Set index to the year.
dax_exec_q.set_index("Sheet1[Year]")

fig = px.pie(dax_exec_q, values='[Total Purchases]', title= 'Purchases per Year', names='Sheet1[Year]', width=1000, height=700)
fig.update_layout(uniformtext_minsize=12, title_font_family="Open Sans", title_font_color="Orange", title_x=0.5, title_xanchor="center", title_font_size=24)
fig.update_traces(textinfo='value', direction='clockwise')
fig.write_image("C:/Users/LHCZ5828/Desktop/Power_BI_training/total_purchases.png")
fig.show()

input("click here to cont.")

#####################################################################################

'''
Create a table and save it as an image.
Attach the image to the PPT file.
'''

# some CSS styling for the table.

th_props = [
        ('font-size', '10px'),
        ('text-align', 'left'),
        ('font-weight', 'bold'),
        ('color', '#021111'), # was #6d6d6d
        ('background-color', '#b7bdbd'),
        ('border-collapse', 'collapse') 
]

td_props = [
    ('border-collapse', 'collapse'),
        ('text-align', 'left')# This doesn't work with pandas, only with CSS.
] 

# Set table styles => a list of dictionaries.
styles = [
dict(selector="th", props=th_props),
dict(selector="td", props=td_props),
dict(selector="", props=[("border", "1px solid black")]) # props => a list of tuples.
]

# P.S; you can rename the columns names to appear differently when displayed as an image.
styled_table = dax_exec_q.style.set_properties(padding="1px", border='2px solid white')\
                                            .set_table_styles(styles)

html = styled_table.render()
options = {
    'width': '100',
    'height': '100',
}
imgkit.from_string(html, "C:/Users/LHCZ5828/Desktop/Power_BI_training/total_purchases_tb.png", options=options)

####################################################################################

'''
Create a PPT file to attach the pie chart.
'''

print("Creating Slides...\n")

img_path = "C:/Users/LHCZ5828/Desktop/Power_BI_training/total_purchases.png"
table_path = "C:/Users/LHCZ5828/Desktop/Power_BI_training/total_purchases_tb.png"

prs = Presentation() # Initialize the Power point class.
blank_slide_layout = prs.slide_layouts[6] # define the type of the slide you want.

slide_0 = prs.slides.add_slide(blank_slide_layout) # Create a slide object with the layout.
slide_1 = prs.slides.add_slide(blank_slide_layout)

top = left =  Inches(1) # To center the figure in the middle of the slide.
height = Inches(5.5)
width = Inches(6)

pic = slide_0.shapes.add_picture(img_path, left, top, height=height) # Attach the pie chart as a static image to the slide.

pic = slide_1.shapes.add_picture(table_path, left, top, height=height) # Attach the table as a static image to the slide.

prs.save('C:/Users/LHCZ5828/Desktop/Power_BI_training/EPPT_test_demo.pptx') # Save the PPT file.

print("\nPPT file saved successfully")
########################################################################################
